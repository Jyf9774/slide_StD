#!/usr/bin/env python3
"""
幻灯片截图分割与PPTX还原数据流水线 v3.2
Slide Screenshot Segmentation and PPTX Reconstruction Pipeline

核心优化：
1. 优先使用Qwen3.6-Plus大模型Grounding能力检测元素，避免传统CV过度分割问题
2. 大模型直接返回元素边界框、类型、文本内容，智能合并相关内容块
3. 传统CV仅作为fallback方案
4. 原图作为背景，元素位置用背景色遮罩后放置提取元素
"""

import os
import json
import hashlib
from pathlib import Path
from dataclasses import dataclass, field
from typing import List, Tuple, Optional, Dict, Any
from enum import Enum
from datetime import datetime

import cv2
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.dml.color import RGBColor

# 导入统一的GPT客户端和Prompt仓库
from openai_client import (
    get_gpt_client,
    encode_image_to_base64,
    encode_numpy_image_to_base64,
    DEFAULT_MODEL
)
import prompts


class GPTAnalyzer:
    """GPT-4.1 视觉分析器"""

    def __init__(self, client=None):
        self.client = client or get_gpt_client()
        self.enabled = self.client is not None

    def analyze_slide(self, image_path: str) -> Dict:
        """
        分析整张幻灯片
        返回：背景色、整体描述、标题等
        """
        if not self.enabled:
            return {}

        base64_image = encode_image_to_base64(image_path)

        try:
            resp = self.client.chat.completions.create(
                model=DEFAULT_MODEL,
                messages=[
                    {"role": "system", "content": prompts.SYSTEM_PROMPT_ANALYZE_SLIDE},
                    {"role": "user", "content": [
                        {"type": "text", "text": prompts.PROMPT_ANALYZE_SLIDE},
                        {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{base64_image}"}}
                    ]}
                ],
                temperature=0.1,
                max_tokens=1000
            )

            result_text = resp.choices[0].message.content.strip()
            # 清理可能的markdown代码块
            if result_text.startswith("```"):
                result_text = result_text.split("```")[1]
                if result_text.startswith("json"):
                    result_text = result_text[4:]
            result_text = result_text.strip()

            return json.loads(result_text)
        except Exception as e:
            print(f"GPT分析幻灯片失败: {e}")
            return {}

    def analyze_element(self, element_image: np.ndarray, context: str = "") -> Dict:
        """
        分析单个元素区域
        返回：类型、文本内容、描述等
        """
        if not self.enabled:
            return {}

        base64_image = encode_numpy_image_to_base64(element_image)

        prompt = prompts.PROMPT_ANALYZE_ELEMENT
        if context:
            prompt += f"\n上下文信息：{context}"

        try:
            resp = self.client.chat.completions.create(
                model=DEFAULT_MODEL,
                messages=[
                    {"role": "system", "content": prompts.SYSTEM_PROMPT_ANALYZE_ELEMENT},
                    {"role": "user", "content": [
                        {"type": "text", "text": prompt},
                        {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{base64_image}"}}
                    ]}
                ],
                temperature=0.1,
                max_tokens=800
            )

            result_text = resp.choices[0].message.content.strip()
            if result_text.startswith("```"):
                result_text = result_text.split("```")[1]
                if result_text.startswith("json"):
                    result_text = result_text[4:]
            result_text = result_text.strip()

            return json.loads(result_text)
        except Exception as e:
            print(f"GPT分析元素失败: {e}")
            return {}
    
    def batch_analyze_elements(self, elements_data: List[Tuple[np.ndarray, str]]) -> List[Dict]:
        """批量分析元素（减少API调用次数）"""
        results = []
        for img, context in elements_data:
            result = self.analyze_element(img, context)
            results.append(result)
        return results

    def detect_elements_grounding(self, image_path: str, image_width: int, image_height: int) -> List[Dict]:
        """
        使用大模型Grounding能力检测整张幻灯片的所有元素
        返回元素列表，每个元素包含：bbox坐标(x, y, width, height)、type、text_content、is_title等信息
        """
        if not self.enabled:
            return []

        base64_image = encode_image_to_base64(image_path)

        prompt = f"""
你是专业的幻灯片元素检测专家，请分析这张幻灯片图片，识别出所有的可视元素。

要求：
1. 识别所有元素：包括文本块、标题、图片、图表、表格、Logo、装饰元素等
2. 每个元素返回精确的边界框坐标：左上角x, 左上角y, 宽度width, 高度height
   - 坐标值必须是整数，基于图片实际尺寸：宽度{image_width}px，高度{image_height}px
   - 边界框要完整包围元素，不要太小或太大
3. 每个元素需要分类：
   - text: 纯文本内容（如果是大标题，is_title设为true）
   - image: 图片、照片、插图
   - chart: 柱状图、折线图、饼图等图表
   - table: 表格
   - diagram: 流程图、架构图等图示
   - logo: 标志、图标
   - decoration: 装饰性元素、分割线等
   - mixed: 混合内容
4. 对于文本元素，提取完整的文本内容
5. 对于其他类型元素，给出简单描述
6. 尽量合并同一语义块的内容，不要过度分割，比如一个完整的段落不要拆分成多个元素
7. 忽略非常小的装饰性元素（小于20x20像素的可以忽略）

返回格式：严格是JSON数组，不需要任何其他内容，不要markdown代码块标记。
每个元素的格式：
{{
  "bbox": [x, y, width, height],
  "type": "元素类型",
  "text_content": "文本内容（文本元素必填，其他类型可选）",
  "is_title": false,
  "description": "元素描述（非文本元素必填）",
  "confidence": 0.95
}}
"""

        try:
            resp = self.client.chat.completions.create(
                model=DEFAULT_MODEL,
                messages=[
                    {"role": "system", "content": "你是专业的幻灯片分析专家，返回严格的JSON格式数据。"},
                    {"role": "user", "content": [
                        {"type": "text", "text": prompt},
                        {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{base64_image}"}}
                    ]}
                ],
                temperature=0.1,
                max_tokens=2000
            )

            result_text = resp.choices[0].message.content.strip()
            # 清理可能的markdown代码块
            if result_text.startswith("```"):
                result_text = result_text.split("```")[1]
                if result_text.startswith("json"):
                    result_text = result_text[4:]
            result_text = result_text.strip()

            import json
            elements = json.loads(result_text)
            return elements
        except Exception as e:
            print(f"大模型Grounding检测失败: {e}")
            return []


# ============ 传统方法备用 ============

class FallbackAnalyzer:
    """传统CV/OCR分析器（GPT不可用时的备用方案）"""
    
    def __init__(self, ocr_lang: str = "chi_sim+eng"):
        self.ocr_lang = ocr_lang
        self._init_tesseract()
    
    def _init_tesseract(self):
        """初始化Tesseract"""
        try:
            import pytesseract
            tesseract_paths = ['/opt/homebrew/bin/tesseract', '/usr/local/bin/tesseract', '/usr/bin/tesseract']
            for path in tesseract_paths:
                if os.path.exists(path):
                    pytesseract.pytesseract.tesseract_cmd = path
                    break
            self.pytesseract = pytesseract
            self.ocr_available = True
        except ImportError:
            self.ocr_available = False
            print("警告: pytesseract未安装，OCR功能不可用")
    
    def extract_background_color(self, image_rgb: np.ndarray) -> str:
        """提取背景色（白色优先策略）"""
        h, w = image_rgb.shape[:2]
        
        # 采样中心区域
        margin_x, margin_y = int(w * 0.15), int(h * 0.15)
        center = image_rgb[margin_y:h-margin_y, margin_x:w-margin_x]
        
        if center.size == 0:
            center = image_rgb
        
        # 白色检测
        white_mask = np.all(center > 240, axis=2)
        white_ratio = np.sum(white_mask) / white_mask.size
        
        if white_ratio > 0.3:
            return "#ffffff"
        
        # 浅色检测
        light_mask = np.all(center > 200, axis=2)
        light_ratio = np.sum(light_mask) / light_mask.size
        
        if light_ratio > 0.4:
            light_pixels = center[light_mask]
            if len(light_pixels) > 0:
                avg = np.mean(light_pixels, axis=0).astype(int)
                return "#{:02x}{:02x}{:02x}".format(*avg)
        
        # K-means聚类
        try:
            from sklearn.cluster import KMeans
            pixels = center.reshape(-1, 3)
            sample = pixels[np.random.choice(len(pixels), min(5000, len(pixels)), replace=False)]
            kmeans = KMeans(n_clusters=3, random_state=42, n_init=10).fit(sample)
            labels, counts = np.unique(kmeans.labels_, return_counts=True)
            dominant = kmeans.cluster_centers_[labels[np.argmax(counts)]].astype(int)
            return "#{:02x}{:02x}{:02x}".format(*dominant)
        except:
            return "#ffffff"
    
    def ocr_region(self, region: np.ndarray) -> str:
        """OCR识别"""
        if not self.ocr_available:
            return ""
        
        try:
            enlarged = cv2.resize(region, None, fx=2, fy=2, interpolation=cv2.INTER_CUBIC)
            pil_image = Image.fromarray(enlarged)
            text = self.pytesseract.image_to_string(pil_image, lang=self.ocr_lang, config='--oem 3 --psm 6')
            return '\n'.join([l.strip() for l in text.split('\n') if l.strip()])
        except Exception as e:
            print(f"OCR错误: {e}")
            return ""
    
    def classify_region(self, region: np.ndarray) -> str:
        """分类区域类型"""
        h, w = region.shape[:2]
        gray = cv2.cvtColor(region, cv2.COLOR_RGB2GRAY)
        
        edges = cv2.Canny(gray, 50, 150)
        edge_density = np.sum(edges > 0) / edges.size
        
        hsv = cv2.cvtColor(region, cv2.COLOR_RGB2HSV)
        h_std, s_std = np.std(hsv[:,:,0]), np.std(hsv[:,:,1])
        
        lines = cv2.HoughLinesP(edges, 1, np.pi/180, 50, minLineLength=30, maxLineGap=10)
        line_count = len(lines) if lines is not None else 0
        
        white_ratio = np.sum(np.all(region > 220, axis=2)) / (h * w)
        
        if w / h > 10 or h / w > 10:
            return "decoration"
        if line_count > 15:
            return "table"
        if line_count > 5 and s_std > 40:
            return "chart"
        if white_ratio > 0.5 and edge_density < 0.15:
            return "text"
        if h_std > 40 and s_std > 50:
            return "image"
        return "mixed"


# ============ 数据结构 ============

class NumpyEncoder(json.JSONEncoder):
    def default(self, obj):
        if isinstance(obj, (np.integer, np.floating, np.bool_)):
            return obj.item()
        if isinstance(obj, np.ndarray):
            return obj.tolist()
        return super().default(obj)


class ElementType(Enum):
    TEXT = "text"
    IMAGE = "image"
    CHART = "chart"
    TABLE = "table"
    DIAGRAM = "diagram"
    LOGO = "logo"
    MIXED = "mixed"
    DECORATION = "decoration"
    UNKNOWN = "unknown"


@dataclass
class BoundingBox:
    x: int
    y: int
    width: int
    height: int
    
    @property
    def x2(self) -> int:
        return self.x + self.width
    
    @property
    def y2(self) -> int:
        return self.y + self.height
    
    @property
    def area(self) -> int:
        return self.width * self.height
    
    def to_dict(self) -> Dict:
        return {"x": self.x, "y": self.y, "width": self.width, "height": self.height}


@dataclass 
class SlideElement:
    id: str
    name: str
    type: ElementType
    bbox: BoundingBox
    image_path: str
    text_content: str = ""
    confidence: float = 0.0
    dominant_colors: List[str] = field(default_factory=list)
    has_border: bool = False
    is_title: bool = False
    z_order: int = 0
    font_size_estimate: int = 12
    text_color: str = "#000000"
    description: str = ""
    metadata: Dict[str, Any] = field(default_factory=dict)
    
    def to_dict(self) -> Dict:
        return {
            "id": self.id,
            "name": self.name,
            "type": self.type.value,
            "bbox": self.bbox.to_dict(),
            "image_path": self.image_path,
            "text_content": self.text_content,
            "confidence": self.confidence,
            "dominant_colors": self.dominant_colors,
            "has_border": self.has_border,
            "is_title": self.is_title,
            "z_order": self.z_order,
            "font_size_estimate": self.font_size_estimate,
            "text_color": self.text_color,
            "description": self.description,
            "metadata": self.metadata
        }


@dataclass
class SlideMetadata:
    slide_id: str
    source_image: str
    width: int
    height: int
    aspect_ratio: str
    background_color: str
    background_colors_palette: List[str]
    element_count: int
    elements: List[SlideElement]
    created_at: str
    version: str = "3.0"
    title: str = ""
    subtitle: str = ""
    slide_type: str = ""
    description: str = ""
    key_points: List[str] = field(default_factory=list)
    keywords: List[str] = field(default_factory=list)
    
    def to_dict(self) -> Dict:
        return {
            "slide_id": self.slide_id,
            "source_image": self.source_image,
            "width": self.width,
            "height": self.height,
            "aspect_ratio": self.aspect_ratio,
            "background_color": self.background_color,
            "background_colors_palette": self.background_colors_palette,
            "element_count": self.element_count,
            "elements": [e.to_dict() for e in self.elements],
            "created_at": self.created_at,
            "version": self.version,
            "title": self.title,
            "subtitle": self.subtitle,
            "slide_type": self.slide_type,
            "description": self.description,
            "key_points": self.key_points,
            "keywords": self.keywords
        }


# ============ 核心流水线 ============

class SlideSegmenter:
    """幻灯片分割器"""

    def __init__(self,
                 min_area: int = 1500,
                 merge_threshold: int = 25,
                 use_gpt: bool = True,
                 detail_level: str = "standard"):
        """
        Args:
            min_area: 最小元素面积，小于这个值的区域会被过滤
            merge_threshold: 合并阈值，小于这个距离的相邻元素会被考虑合并
            use_gpt: 是否使用GPT分析
            detail_level: 分割精细度: "fine"(精细)/"standard"(标准)/"coarse"(粗略)
        """
        self.detail_level = detail_level.lower()

        # 根据精细度调整基础参数
        if self.detail_level == "fine":
            self.min_area = min_area or 200
            self.merge_threshold = merge_threshold or 25
            self.dilate_iterations = 1
            self.canny_low = 20
            self.canny_high = 100
        elif self.detail_level == "coarse":
            self.min_area = min_area or 500
            self.merge_threshold = merge_threshold or 50
            self.dilate_iterations = 2
            self.canny_low = 30
            self.canny_high = 120
        else: # standard
            self.min_area = min_area or 300
            self.merge_threshold = merge_threshold or 35
            self.dilate_iterations = 1
            self.canny_low = 25
            self.canny_high = 110

        self.use_gpt = use_gpt

        # 初始化分析器
        self.gpt_analyzer = GPTAnalyzer() if use_gpt else None
        self.fallback_analyzer = FallbackAnalyzer()

        # 判断是否可以使用GPT
        self.gpt_available = self.gpt_analyzer and self.gpt_analyzer.enabled
        if use_gpt and not self.gpt_available:
            print("提示: GPT-4.1不可用，将使用传统CV/OCR方法")
        
    def segment(self, image_path: str, output_dir: str) -> SlideMetadata:
        """分割幻灯片图像"""
        output_path = Path(output_dir)
        output_path.mkdir(parents=True, exist_ok=True)
        elements_dir = output_path / "elements"
        elements_dir.mkdir(exist_ok=True)
        
        # 读取图像
        image = cv2.imread(image_path)
        if image is None:
            raise ValueError(f"无法读取图像: {image_path}")
        
        height, width = image.shape[:2]
        image_rgb = cv2.cvtColor(image, cv2.COLOR_BGR2RGB)
        
        slide_id = self._generate_id(image_path)
        
        # ====== 使用GPT分析整体幻灯片 ======
        slide_analysis = {}
        if self.gpt_available:
            print("正在使用GPT-4.1分析幻灯片...")
            slide_analysis = self.gpt_analyzer.analyze_slide(image_path)
        
        # 提取背景色
        if slide_analysis.get('background_color'):
            bg_color = slide_analysis['background_color']
        else:
            bg_color = self.fallback_analyzer.extract_background_color(image_rgb)
        
        bg_palette = [bg_color]
        
        # ====== 优先使用大模型Grounding检测元素 ======
        raw_elements = []
        if self.use_gpt and self.gpt_available:
            print("使用大模型Grounding检测元素...")
            grounding_elements = self._gpt_grounding_detection(image_path, image_rgb, str(elements_dir), slide_id)
            if grounding_elements:
                raw_elements = grounding_elements
                print(f"大模型检测到 {len(raw_elements)} 个元素")

        # ====== Fallback: 使用传统CV检测 ======
        if not raw_elements:
            print("使用传统CV检测元素...")
            regions = self._detect_regions(image)
            print(f"CV检测到 {len(regions)} 个候选区域")

            for idx, bbox in enumerate(regions):
                element = self._process_region(
                    image_rgb, bbox, idx,
                    str(elements_dir), slide_id,
                    slide_analysis.get('title', '')
                )
                if element:
                    raw_elements.append(element)

            # ====== 智能合并同类型相邻区域 ======
            elements = self._smart_merge_elements(raw_elements)
            print(f"智能合并后剩余 {len(elements)} 个元素")

            # ====== 后处理：过滤和优化元素 ======
            elements = self._post_process_elements(elements)
            print(f"后处理后剩余 {len(elements)} 个元素")
        else:
            elements = raw_elements

        # 排序
        elements.sort(key=lambda e: (e.bbox.y, e.bbox.x))
        for idx, elem in enumerate(elements):
            elem.z_order = idx
        
        # 创建元数据
        metadata = SlideMetadata(
            slide_id=slide_id,
            source_image=os.path.basename(image_path),
            width=width,
            height=height,
            aspect_ratio=self._calc_aspect_ratio(width, height),
            background_color=bg_color,
            background_colors_palette=bg_palette,
            element_count=len(elements),
            elements=elements,
            created_at=datetime.now().isoformat(),
            title=slide_analysis.get('title', ''),
            subtitle=slide_analysis.get('subtitle', ''),
            slide_type=slide_analysis.get('slide_type', ''),
            description=slide_analysis.get('description', ''),
            key_points=slide_analysis.get('key_points', [])
        )
        
        # 提取关键词
        all_text = " ".join([e.text_content for e in elements if e.text_content])
        metadata.keywords = list(set([w for w in all_text.split() if len(w) > 2]))[:30]
        
        # 保存JSON
        json_path = output_path / f"{slide_id}.json"
        with open(json_path, 'w', encoding='utf-8') as f:
            json.dump(metadata.to_dict(), f, ensure_ascii=False, indent=2, cls=NumpyEncoder)
        
        # 复制原图
        original_copy = output_path / f"original_{os.path.basename(image_path)}"
        cv2.imwrite(str(original_copy), image)
        
        return metadata
    
    def _generate_id(self, image_path: str) -> str:
        content = f"{image_path}_{datetime.now().isoformat()}"
        return hashlib.md5(content.encode()).hexdigest()[:12]
    
    def _detect_regions(self, image: np.ndarray) -> List[BoundingBox]:
        """CV检测区域边界框（优化版，减少细碎区域）"""
        gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
        h, w = image.shape[:2]

        # 自适应阈值（提高灵敏度）
        binary = cv2.adaptiveThreshold(gray, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
                                        cv2.THRESH_BINARY_INV, 11, 2)

        # Canny边缘检测（使用动态阈值，降低阈值提高灵敏度）
        edges = cv2.Canny(gray, self.canny_low, self.canny_high)
        combined = cv2.bitwise_or(binary, edges)

        # 形态学操作（调整kernel大小，平衡合并效果和区域完整性）
        kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (12, 4))
        dilated = cv2.dilate(combined, kernel, iterations=self.dilate_iterations)
        closed = cv2.morphologyEx(dilated, cv2.MORPH_CLOSE,
                                   cv2.getStructuringElement(cv2.MORPH_RECT, (6, 6)))

        contours, _ = cv2.findContours(closed, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)

        regions = []
        for contour in contours:
            x, y, cw, ch = cv2.boundingRect(contour)
            area = cw * ch

            if area < self.min_area or area > 0.98 * w * h:
                continue

            aspect = cw / ch if ch > 0 else 0
            # 稍微严格的宽高比过滤，去掉极端细长的元素
            if aspect > 25 or aspect < 0.04:
                continue

            padding = 8  # 增加padding，让边界框包含更多上下文，方便后续合并
            x = max(0, x - padding)
            y = max(0, y - padding)
            cw = min(w - x, cw + 2 * padding)
            ch = min(h - y, ch + 2 * padding)

            regions.append(BoundingBox(x, y, cw, ch))

        # 先进行一次简单的重叠合并，去掉完全包含在其他区域内的小区域
        regions = self._remove_nested_regions(regions)

        return regions

    def _remove_nested_regions(self, regions: List[BoundingBox]) -> List[BoundingBox]:
        """移除完全被其他区域包含的嵌套小区域"""
        if len(regions) <= 1:
            return regions

        # 按面积从大到小排序
        sorted_regions = sorted(regions, key=lambda r: -r.area)
        kept = []

        for region in sorted_regions:
            # 检查当前区域是否被已经保留的区域完全包含
            nested = False
            for kept_region in kept:
                if (kept_region.x <= region.x and
                    kept_region.y <= region.y and
                    kept_region.x2 >= region.x2 and
                    kept_region.y2 >= region.y2):
                    nested = True
                    break

            if not nested:
                kept.append(region)

        # 恢复原来的顺序（按位置排序）
        kept.sort(key=lambda r: (r.y, r.x))
        return kept
    
    def _merge_overlapping_regions(self, regions: List[BoundingBox]) -> List[BoundingBox]:
        """旧的粗暴合并方法，已废弃，保留用于兼容"""
        if not regions:
            return []

        boxes = [[r.x, r.y, r.x2, r.y2] for r in regions]

        merged = True
        while merged:
            merged = False
            new_boxes = []
            used = set()

            for i, box1 in enumerate(boxes):
                if i in used:
                    continue
                x1, y1, x2, y2 = box1

                for j, box2 in enumerate(boxes[i+1:], i+1):
                    if j in used:
                        continue
                    bx1, by1, bx2, by2 = box2

                    if (x1 - self.merge_threshold <= bx2 and
                        x2 + self.merge_threshold >= bx1 and
                        y1 - self.merge_threshold <= by2 and
                        y2 + self.merge_threshold >= by1):
                        x1, y1 = min(x1, bx1), min(y1, by1)
                        x2, y2 = max(x2, bx2), max(y2, by2)
                        used.add(j)
                        merged = True

                new_boxes.append([x1, y1, x2, y2])
                used.add(i)

            boxes = new_boxes

        return [BoundingBox(b[0], b[1], b[2]-b[0], b[3]-b[1]) for b in boxes]

    def _smart_merge_elements(self, elements: List[SlideElement]) -> List[SlideElement]:
        """
        智能合并元素（增强版）：
        1. 先过滤掉无用的小装饰元素
        2. 图片/图表/表格/Logo等实体元素不合并
        3. 文本元素智能合并：同一行文本自动合并，相邻行的段落自动合并
        4. 混合类型元素如果内容相关也可以合并
        """
        if len(elements) <= 1:
            return elements

        # 第一步：先过滤掉太小的装饰性元素
        filtered_elements = []
        min_area_for_decoration = self.min_area * 2  # 装饰元素需要至少是最小面积的2倍才保留
        for elem in elements:
            # 装饰元素特别小的直接过滤
            if elem.type == ElementType.DECORATION and elem.bbox.area < min_area_for_decoration:
                continue
            # 其他类型元素如果面积特别小（小于最小面积的一半）也过滤
            if elem.bbox.area < self.min_area * 0.5:
                continue
            filtered_elements.append(elem)

        if len(filtered_elements) <= 1:
            return filtered_elements

        # 按y坐标排序，从上到下处理
        elements = sorted(filtered_elements, key=lambda e: (e.bbox.y, e.bbox.x))

        merged = []
        used = set()

        for i, elem1 in enumerate(elements):
            if i in used:
                continue

            # 图片/图表/表格/Logo/图示类型不合并，直接保留
            if elem1.type in [ElementType.IMAGE, ElementType.CHART, ElementType.TABLE,
                              ElementType.LOGO, ElementType.DIAGRAM]:
                merged.append(elem1)
                used.add(i)
                continue

            current_elem = elem1
            used.add(i)

            # 查找可以合并的相邻元素
            for j, elem2 in enumerate(elements[i+1:], i+1):
                if j in used:
                    continue

                bbox1 = current_elem.bbox
                bbox2 = elem2.bbox

                # 检查是否距离太远
                horizontal_dist = max(0, max(bbox1.x, bbox2.x) - min(bbox1.x2, bbox2.x2))
                vertical_dist = max(0, max(bbox1.y, bbox2.y) - min(bbox1.y2, bbox2.y2))

                if (horizontal_dist > self.merge_threshold * 2 or
                    vertical_dist > self.merge_threshold * 3):
                    continue

                # 类型兼容性检查
                type_compatible = False
                # 完全相同类型
                if elem2.type == current_elem.type:
                    type_compatible = True
                # 文本和混合类型可以互相合并
                elif {current_elem.type, elem2.type} <= {ElementType.TEXT, ElementType.MIXED}:
                    type_compatible = True
                # 两个都是装饰类型也可以合并
                elif current_elem.type == ElementType.DECORATION and elem2.type == ElementType.DECORATION:
                    type_compatible = True

                if not type_compatible:
                    continue

                # 文本/混合类型额外判断
                if current_elem.type in [ElementType.TEXT, ElementType.MIXED]:
                    height1 = bbox1.height
                    height2 = bbox2.height
                    max_height = max(height1, height2)
                    min_height = min(height1, height2)

                    # 高度差不能太大
                    if max_height > min_height * 2.5:
                        continue

                    # 判断是水平相邻（同一行）还是垂直相邻（上下行）
                    y_center1 = bbox1.y + height1 / 2
                    y_center2 = bbox2.y + height2 / 2
                    y_center_diff = abs(y_center1 - y_center2)

                    # 同一行文本（y中心差小于行高的50%）
                    if y_center_diff < max_height * 0.5:
                        # 同一行的文本，只要水平距离不太大都可以合并
                        if horizontal_dist > self.merge_threshold * 3:
                            continue
                    # 上下行文本（垂直距离小于行高的2倍）
                    elif vertical_dist < max_height * 2.0:
                        # 上下行的文本，水平方向需要有一定重叠
                        x_overlap = min(bbox1.x2, bbox2.x2) - max(bbox1.x, bbox2.x)
                        min_width = min(bbox1.width, bbox2.width)
                        if x_overlap < -min_width * 0.2:  # 允许稍微错开20%
                            continue
                    else:
                        # 既不在同一行也不是相邻行，不合并
                        continue

                # 可以合并，合并两个元素
                new_bbox = BoundingBox(
                    x=min(bbox1.x, bbox2.x),
                    y=min(bbox1.y, bbox2.y),
                    width=max(bbox1.x2, bbox2.x2) - min(bbox1.x, bbox2.x),
                    height=max(bbox1.y2, bbox2.y2) - min(bbox1.y, bbox2.y)
                )

                # 合并文本内容，用换行分隔
                combined_text = []
                if current_elem.text_content:
                    combined_text.append(current_elem.text_content)
                if elem2.text_content:
                    combined_text.append(elem2.text_content)
                combined_text = "\n".join(combined_text).strip()

                # 合并描述
                combined_description = []
                if current_elem.description:
                    combined_description.append(current_elem.description)
                if elem2.description:
                    combined_description.append(elem2.description)
                combined_description = " ".join(combined_description).strip()

                # 更新当前元素的属性
                current_elem.bbox = new_bbox
                current_elem.text_content = combined_text
                current_elem.description = combined_description

                # 类型如果一个是TEXT一个是MIXED，取TEXT
                if current_elem.type == ElementType.MIXED and elem2.type == ElementType.TEXT:
                    current_elem.type = ElementType.TEXT
                elif current_elem.type == ElementType.TEXT and elem2.type == ElementType.MIXED:
                    current_elem.type = ElementType.TEXT  # 保持TEXT类型

                # 保留更可能是标题的属性
                current_elem.is_title = current_elem.is_title or elem2.is_title

                # 字体大小取较大的那个（通常标题字体更大）
                current_elem.font_size_estimate = max(current_elem.font_size_estimate,
                                                     elem2.font_size_estimate)

                # 置信度取平均值
                current_elem.confidence = (current_elem.confidence + elem2.confidence) / 2

                # 合并主要颜色
                current_elem.dominant_colors = list(set(current_elem.dominant_colors + elem2.dominant_colors))[:3]

                used.add(j)

            merged.append(current_elem)

        # 二次过滤：合并后再次检查有没有太小的元素
        final_merged = []
        for elem in merged:
            if elem.bbox.area >= self.min_area * 0.7:  # 允许稍微小于最小面积
                final_merged.append(elem)

        # 重新排序和分配z_order
        final_merged.sort(key=lambda e: (e.bbox.y, e.bbox.x))
        for idx, elem in enumerate(final_merged):
            elem.z_order = idx

        return final_merged
    
    def _process_region(self, image_rgb: np.ndarray, bbox: BoundingBox, 
                        idx: int, output_dir: str, slide_id: str,
                        slide_title: str = "") -> Optional[SlideElement]:
        """处理单个区域"""
        crop = image_rgb[bbox.y:bbox.y2, bbox.x:bbox.x2]
        
        if crop.size == 0:
            return None
        
        # 先进行分析以获取类型信息（用于命名）
        if self.gpt_available:
            context = f"这是幻灯片'{slide_title}'中的一个区域" if slide_title else ""
            analysis = self.gpt_analyzer.analyze_element(crop, context)
            elem_type_str = analysis.get('type', 'mixed')
            text_content = analysis.get('text_content', '')
            description = analysis.get('description', '')
            is_title = analysis.get('is_title', False)
            font_size = analysis.get('estimated_font_size', 12)
            text_color = analysis.get('text_color', '#000000')
            has_border = analysis.get('has_border', False)
            confidence = 0.9
        else:
            elem_type_str = self.fallback_analyzer.classify_region(crop)
            text_content = self.fallback_analyzer.ocr_region(crop) if elem_type_str in ['text', 'mixed', 'table'] else ""
            description = ""
            is_title = bbox.y < image_rgb.shape[0] * 0.25 and bbox.width > image_rgb.shape[1] * 0.3
            font_size = int(bbox.height / max(1, len(text_content.split('\n'))) * 0.6) if text_content else 12
            font_size = max(8, min(72, font_size))
            text_color = "#000000"
            has_border = self._detect_border(crop)
            confidence = 0.6
        
        # 生成有意义的元素名称
        # 格式: {序号}_{类型}_{简短标识}
        element_name = self._generate_element_name(idx, elem_type_str, text_content, is_title)
        element_id = f"{slide_id}_{element_name}"
        
        # 保存裁剪图像（使用新名称）
        image_filename = f"{element_name}.png"
        image_path = os.path.join(output_dir, image_filename)
        Image.fromarray(crop).save(image_path)
        
        # 转换类型枚举
        try:
            elem_type = ElementType(elem_type_str.lower())
        except ValueError:
            elem_type = ElementType.MIXED
        
        # 提取主要颜色
        dominant_colors = self._get_dominant_colors(crop)
        
        return SlideElement(
            id=element_id,
            name=element_name,
            type=elem_type,
            bbox=bbox,
            image_path=f"elements/{image_filename}",
            text_content=text_content,
            confidence=confidence,
            dominant_colors=dominant_colors,
            has_border=has_border,
            is_title=is_title,
            font_size_estimate=font_size,
            text_color=text_color,
            description=description,
            metadata={
                "aspect_ratio": round(bbox.width / bbox.height, 2) if bbox.height > 0 else 0,
                "relative_size": bbox.area
            }
        )
    
    def _generate_element_name(self, idx: int, elem_type: str, text_content: str, is_title: bool) -> str:
        """
        生成有意义的元素名称
        格式: {序号:02d}_{类型}_{简短标识}
        
        示例:
        - 00_text_title
        - 01_text_abstract
        - 02_table_content
        - 03_logo_alibaba
        - 04_chart_data
        """
        import re
        
        # 基础前缀
        prefix = f"{idx:02d}_{elem_type}"
        
        # 生成简短标识
        if is_title:
            suffix = "title"
        elif text_content:
            # 从文本内容提取关键词
            # 清理文本
            clean_text = text_content.strip().split('\n')[0]  # 取第一行
            clean_text = re.sub(r'[^\w\s]', '', clean_text)  # 移除标点
            words = clean_text.split()
            
            if words:
                # 取前1-2个有意义的词
                key_words = []
                for w in words[:3]:
                    if len(w) > 1:  # 跳过单字符
                        # 如果是中文，取前4个字符
                        if any('\u4e00' <= c <= '\u9fff' for c in w):
                            key_words.append(w[:4])
                        else:
                            key_words.append(w[:10].lower())
                    if len(key_words) >= 2:
                        break
                
                suffix = "_".join(key_words) if key_words else "content"
            else:
                suffix = "content"
        else:
            # 无文本内容，使用类型相关的默认名称
            type_defaults = {
                "image": "img",
                "chart": "chart",
                "table": "table",
                "logo": "logo",
                "diagram": "diagram",
                "decoration": "deco",
                "mixed": "mixed"
            }
            suffix = type_defaults.get(elem_type, "elem")
        
        # 清理名称，确保是有效的标识符
        suffix = re.sub(r'[^\w]', '_', suffix)
        suffix = re.sub(r'_+', '_', suffix).strip('_')
        
        if not suffix:
            suffix = "elem"
        
        return f"{prefix}_{suffix}"
    
    def _get_dominant_colors(self, region: np.ndarray, n_colors: int = 3) -> List[str]:
        """提取主要颜色"""
        try:
            from sklearn.cluster import KMeans
            pixels = region.reshape(-1, 3)
            sample = pixels[np.random.choice(len(pixels), min(3000, len(pixels)), replace=False)]
            kmeans = KMeans(n_clusters=min(n_colors, len(sample)//10+1), random_state=42, n_init=10).fit(sample)
            return ["#{:02x}{:02x}{:02x}".format(*c.astype(int)) for c in kmeans.cluster_centers_]
        except:
            return []
    
    def _detect_border(self, region: np.ndarray) -> bool:
        """检测边框"""
        gray = cv2.cvtColor(region, cv2.COLOR_RGB2GRAY)
        edges = cv2.Canny(gray, 50, 150)
        h, w = edges.shape
        bw = 3
        
        scores = [
            np.sum(edges[:bw, :] > 0) / (bw * w) if w > 0 else 0,
            np.sum(edges[-bw:, :] > 0) / (bw * w) if w > 0 else 0,
            np.sum(edges[:, :bw] > 0) / (h * bw) if h > 0 else 0,
            np.sum(edges[:, -bw:] > 0) / (h * bw) if h > 0 else 0
        ]
        return sum(s > 0.3 for s in scores) >= 3
    
    def _post_process_elements(self, elements: List[SlideElement]) -> List[SlideElement]:
        """
        后处理元素：
        1. 过滤掉完全没有内容的元素
        2. 合并高度重叠的元素
        3. 调整元素属性
        """
        if len(elements) <= 1:
            return elements

        # 第一步：过滤没有内容的元素
        filtered = []
        for elem in elements:
            # 如果是文本/混合类型，完全没有文本内容且面积很小，过滤
            if elem.type in [ElementType.TEXT, ElementType.MIXED]:
                if not elem.text_content and elem.bbox.area < self.min_area * 1.5:
                    continue
            # 装饰元素没有边框且面积很小，过滤
            if elem.type == ElementType.DECORATION and not elem.has_border and elem.bbox.area < self.min_area * 2:
                continue
            filtered.append(elem)

        if len(filtered) <= 1:
            return filtered

        # 第二步：合并高度重叠的元素
        # 按面积从大到小排序
        sorted_by_area = sorted(filtered, key=lambda e: -e.bbox.area)
        kept = []
        used = set()

        for i, elem1 in enumerate(sorted_by_area):
            if i in used:
                continue

            current = elem1
            used.add(i)

            for j, elem2 in enumerate(sorted_by_area[i+1:], i+1):
                if j in used:
                    continue

                bbox1 = current.bbox
                bbox2 = elem2.bbox

                # 计算重叠面积
                x_overlap = max(0, min(bbox1.x2, bbox2.x2) - max(bbox1.x, bbox2.x))
                y_overlap = max(0, min(bbox1.y2, bbox2.y2) - max(bbox1.y, bbox2.y))
                overlap_area = x_overlap * y_overlap

                # 如果重叠面积超过小元素面积的70%，合并
                min_area = min(bbox1.area, bbox2.area)
                if overlap_area > min_area * 0.7:
                    # 合并到较大的元素
                    new_bbox = BoundingBox(
                        x=min(bbox1.x, bbox2.x),
                        y=min(bbox1.y, bbox2.y),
                        width=max(bbox1.x2, bbox2.x2) - min(bbox1.x, bbox2.x),
                        height=max(bbox1.y2, bbox2.y2) - min(bbox1.y, bbox2.y)
                    )
                    current.bbox = new_bbox

                    # 合并文本内容
                    if elem2.text_content and elem2.text_content not in current.text_content:
                        current.text_content = f"{current.text_content}\n{elem2.text_content}".strip()

                    # 合并描述
                    if elem2.description and elem2.description not in current.description:
                        current.description = f"{current.description} {elem2.description}".strip()

                    # 保留更可能是标题的属性
                    current.is_title = current.is_title or elem2.is_title

                    used.add(j)

            kept.append(current)

        # 恢复按位置排序
        kept.sort(key=lambda e: (e.bbox.y, e.bbox.x))
        return kept

    def _gpt_grounding_detection(self, image_path: str, image_rgb: np.ndarray, elements_dir: str, slide_id: str) -> List[SlideElement]:
        """使用大模型Grounding检测元素，直接返回SlideElement列表"""
        height, width = image_rgb.shape[:2]

        # 调用大模型Grounding检测
        elements_data = self.gpt_analyzer.detect_elements_grounding(image_path, width, height)
        if not elements_data:
            return []

        elements = []
        for idx, elem_data in enumerate(elements_data):
            try:
                # 解析边界框
                bbox_data = elem_data.get("bbox", [0, 0, 0, 0])
                x, y, w, h = map(int, bbox_data)

                # 坐标合法性检查
                x = max(0, x)
                y = max(0, y)
                w = min(width - x, w)
                h = min(height - y, h)

                if w <= 0 or h <= 0 or w * h < self.min_area:
                    continue

                bbox = BoundingBox(x, y, w, h)

                # 裁剪元素图像
                crop = image_rgb[y:y+h, x:x+w]
                if crop.size == 0:
                    continue

                # 元素类型
                elem_type_str = elem_data.get("type", "mixed").lower()
                try:
                    elem_type = ElementType(elem_type_str)
                except ValueError:
                    elem_type = ElementType.MIXED

                # 生成元素名称
                text_content = elem_data.get("text_content", "")
                is_title = elem_data.get("is_title", False)
                element_name = self._generate_element_name(idx, elem_type_str, text_content, is_title)
                element_id = f"{slide_id}_{element_name}"

                # 保存元素图像
                image_filename = f"{element_name}.png"
                image_path = os.path.join(elements_dir, image_filename)
                Image.fromarray(crop).save(image_path)

                # 创建SlideElement
                element = SlideElement(
                    id=element_id,
                    name=element_name,
                    type=elem_type,
                    bbox=bbox,
                    image_path=f"elements/{image_filename}",
                    text_content=text_content,
                    confidence=elem_data.get("confidence", 0.9),
                    is_title=is_title,
                    description=elem_data.get("description", ""),
                    font_size_estimate=self._estimate_font_size(crop, text_content)
                )

                elements.append(element)
            except Exception as e:
                print(f"处理大模型检测到的元素失败: {e}")
                continue

        return elements

    def _estimate_font_size(self, crop: np.ndarray, text_content: str) -> int:
        """估算文本元素的字体大小"""
        if not text_content:
            return 12
        h = crop.shape[0]
        lines = len(text_content.split("\n"))
        if lines == 0:
            lines = 1
        return max(8, min(72, int(h / lines * 0.6)))

    def _calc_aspect_ratio(self, width: int, height: int) -> str:
        from math import gcd
        d = gcd(width, height)
        return f"{width//d}:{height//d}"


class SlideReconstructor:
    """幻灯片重建器 - 原图背景 + 元素遮罩方案"""
    
    def __init__(self, 
                 slide_width_inches: float = 13.333, 
                 slide_height_inches: float = 7.5,
                 use_original_as_background: bool = True,
                 mask_elements: bool = True):
        """
        Args:
            use_original_as_background: 使用原图作为背景
            mask_elements: 在元素位置用背景色遮罩
        """
        self.slide_width = Inches(slide_width_inches)
        self.slide_height = Inches(slide_height_inches)
        self.use_original_as_background = use_original_as_background
        self.mask_elements = mask_elements
    
    def reconstruct(self, metadata_path: str, output_path: str) -> str:
        """从元数据重建幻灯片"""
        with open(metadata_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        
        metadata_dir = os.path.dirname(metadata_path)
        
        prs = Presentation()
        prs.slide_width = self.slide_width
        prs.slide_height = self.slide_height
        
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        # 计算缩放比例
        original_width = data['width']
        original_height = data['height']
        scale_x = self.slide_width / original_width
        scale_y = self.slide_height / original_height
        
        bg_color = data.get('background_color', '#ffffff')
        
        # ====== 方案：原图背景 + 遮罩 ======
        if self.use_original_as_background:
            # 1. 查找原图
            original_image_path = None
            for fname in os.listdir(metadata_dir):
                if fname.startswith("original_"):
                    original_image_path = os.path.join(metadata_dir, fname)
                    break
            
            if original_image_path and os.path.exists(original_image_path):
                if self.mask_elements:
                    # 创建带遮罩的背景图
                    masked_bg_path = self._create_masked_background(
                        original_image_path, data['elements'], bg_color, metadata_dir
                    )
                    # 插入遮罩后的背景图
                    slide.shapes.add_picture(masked_bg_path, 0, 0, self.slide_width, self.slide_height)
                else:
                    # 直接插入原图作为背景
                    slide.shapes.add_picture(original_image_path, 0, 0, self.slide_width, self.slide_height)
            else:
                # 没有原图，使用纯色背景
                self._set_background_color(slide, bg_color)
        else:
            self._set_background_color(slide, bg_color)
        
        # ====== 插入元素 ======
        elements = sorted(data['elements'], key=lambda e: e.get('z_order', 0))
        
        for elem_data in elements:
            self._add_element(slide, elem_data, metadata_dir, scale_x, scale_y)
        
        prs.save(output_path)
        return output_path
    
    def _create_masked_background(self, original_path: str, elements: List[Dict], 
                                   bg_color: str, output_dir: str) -> str:
        """创建带遮罩的背景图"""
        # 读取原图
        image = cv2.imread(original_path)
        image_rgb = cv2.cvtColor(image, cv2.COLOR_BGR2RGB)
        
        # 解析背景色
        bg_hex = bg_color.lstrip('#')
        bg_rgb = (int(bg_hex[0:2], 16), int(bg_hex[2:4], 16), int(bg_hex[4:6], 16))
        
        # 在元素位置绘制背景色矩形
        for elem in elements:
            bbox = elem['bbox']
            x, y = bbox['x'], bbox['y']
            w, h = bbox['width'], bbox['height']
            
            # 稍微扩大遮罩区域，确保完全覆盖
            padding = 2
            x1 = max(0, x - padding)
            y1 = max(0, y - padding)
            x2 = min(image_rgb.shape[1], x + w + padding)
            y2 = min(image_rgb.shape[0], y + h + padding)
            
            # 填充背景色
            image_rgb[y1:y2, x1:x2] = bg_rgb
        
        # 保存遮罩后的图像
        masked_path = os.path.join(output_dir, "masked_background.png")
        Image.fromarray(image_rgb).save(masked_path)
        
        return masked_path
    
    def _set_background_color(self, slide, hex_color: str):
        """设置幻灯片背景色"""
        hex_color = hex_color.lstrip('#')
        r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(r, g, b)
    
    def _add_element(self, slide, elem_data: Dict, base_dir: str, 
                     scale_x: Emu, scale_y: Emu):
        """添加元素"""
        bbox = elem_data['bbox']
        
        left = int(bbox['x'] * scale_x)
        top = int(bbox['y'] * scale_y)
        width = int(bbox['width'] * scale_x)
        height = int(bbox['height'] * scale_y)
        
        image_path = os.path.join(base_dir, elem_data['image_path'])
        
        if os.path.exists(image_path):
            shape = slide.shapes.add_picture(image_path, left, top, width, height)
            shape.name = elem_data['name']
        else:
            print(f"警告: 图片不存在 {image_path}")


# ============ 便捷函数 ============

def process_slide(image_path: str,
                  output_base_dir: str = "./output",
                  use_gpt: bool = True,
                  use_original_bg: bool = True,
                  mask_elements: bool = True,
                  detail_level: str = "standard",
                  min_area: int = None,
                  merge_threshold: int = None) -> Tuple[str, str]:
    """
    处理单张幻灯片

    Args:
        image_path: 输入图像路径
        output_base_dir: 输出基础目录
        use_gpt: 是否使用GPT-4.1进行分析
        use_original_bg: 是否使用原图作为背景
        mask_elements: 是否在元素位置遮罩
        detail_level: 分割精细度: "fine"(精细)/"standard"(标准)/"coarse"(粗略)
        min_area: 自定义最小元素面积（覆盖精细度的默认值）
        merge_threshold: 自定义合并阈值（覆盖精细度的默认值）

    Returns:
        (json_path, pptx_path): JSON和PPTX文件路径
    """
    image_name = os.path.splitext(os.path.basename(image_path))[0]
    output_dir = os.path.join(output_base_dir, image_name)

    # 分割
    segmenter = SlideSegmenter(
        use_gpt=use_gpt,
        detail_level=detail_level,
        min_area=min_area,
        merge_threshold=merge_threshold
    )
    metadata = segmenter.segment(image_path, output_dir)
    
    json_path = os.path.join(output_dir, f"{metadata.slide_id}.json")
    
    # 重建
    pptx_path = os.path.join(output_dir, f"{image_name}_reconstructed.pptx")
    reconstructor = SlideReconstructor(
        use_original_as_background=use_original_bg,
        mask_elements=mask_elements
    )
    reconstructor.reconstruct(json_path, pptx_path)
    
    return json_path, pptx_path


if __name__ == "__main__":
    import sys
    
    print("""
幻灯片分割与重建流水线 v3.1
===========================
用法: python pipeline.py <image_path> [output_dir] [options]

选项:
  --no-gpt        不使用GPT-4.1（使用传统CV/OCR）
  --no-bg         不使用原图作为背景
  --no-mask       不遮罩元素区域
  --fine          精细分割模式（检测更多小元素）
  --coarse        粗略分割模式（合并更多元素，减少分割数量）
  --min-area N    自定义最小元素面积（默认根据精细度自动调整）
  --merge-thresh N 自定义合并阈值（默认根据精细度自动调整）

示例:
  python pipeline.py slide.png ./output
  python pipeline.py slide.png ./output --no-gpt
  python pipeline.py slide.png ./output --coarse  # 粗略分割，减少元素数量
  python pipeline.py slide.png ./output --fine    # 精细分割，保留更多细节
""")

    if len(sys.argv) < 2:
        sys.exit(1)

    image_path = sys.argv[1]
    output_dir = "./output"
    use_gpt = True
    use_original_bg = True
    mask_elements = True
    detail_level = "standard"
    min_area = None
    merge_threshold = None

    i = 2
    while i < len(sys.argv):
        arg = sys.argv[i]
        if arg == "--no-gpt":
            use_gpt = False
        elif arg == "--no-bg":
            use_original_bg = False
        elif arg == "--no-mask":
            mask_elements = False
        elif arg == "--fine":
            detail_level = "fine"
        elif arg == "--coarse":
            detail_level = "coarse"
        elif arg == "--min-area" and i + 1 < len(sys.argv):
            try:
                min_area = int(sys.argv[i+1])
                i += 1
            except ValueError:
                print("警告: --min-area 参数需要是整数，使用默认值")
        elif arg == "--merge-thresh" and i + 1 < len(sys.argv):
            try:
                merge_threshold = int(sys.argv[i+1])
                i += 1
            except ValueError:
                print("警告: --merge-thresh 参数需要是整数，使用默认值")
        elif not arg.startswith("-"):
            output_dir = arg
        i += 1
    
    print(f"处理图像: {image_path}")
    print(f"输出目录: {output_dir}")
    print(f"使用GPT: {use_gpt}")
    print(f"原图背景: {use_original_bg}")
    print(f"元素遮罩: {mask_elements}")
    print(f"分割精细度: {detail_level}")
    if min_area is not None:
        print(f"自定义最小面积: {min_area}")
    if merge_threshold is not None:
        print(f"自定义合并阈值: {merge_threshold}")
    print()

    json_path, pptx_path = process_slide(
        image_path, output_dir,
        use_gpt=use_gpt,
        use_original_bg=use_original_bg,
        mask_elements=mask_elements,
        detail_level=detail_level,
        min_area=min_area,
        merge_threshold=merge_threshold
    )
    
    print(f"\n处理完成!")
    print(f"JSON: {json_path}")
    print(f"PPTX: {pptx_path}")
