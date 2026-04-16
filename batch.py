#!/usr/bin/env python3
"""
幻灯片处理流水线 - 命令行工具
完整工作流程：幻灯片截图 → 元素分割 → PPTX还原 → 旁白生成 → TTS语音合成 → 智能动画生成
"""

import os
import json
import glob
import argparse
from pathlib import Path
from dataclasses import dataclass, asdict
from typing import List, Dict, Optional, Callable
from concurrent.futures import ThreadPoolExecutor, as_completed
import logging

# PPTX相关导入
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.dml.color import RGBColor

# 内部模块导入（统一使用新版v4流水线）
from pipeline import SlideProcessor, SlideMetadata, process_slide
from narration_generator import generate_narration, generate_batch_narrations
from media import generate_tts_and_animations

# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)


@dataclass
class PipelineConfig:
    """流水线配置"""
    # 分割参数
    min_area: int = 300  # 最小语义块面积，像素，小于这个值的块会被过滤
    ocr_lang: str = "chi_sim+eng"
    hybrid_mode: bool = True  # DocLayout-YOLO版面检测 + VLM语义增强的混合模式

    # 处理参数
    max_workers: int = 4
    save_intermediate: bool = True
    verbose: bool = True  # 是否输出详细调试信息

    # 输出控制
    generate_description: bool = True  # 生成描述（为TTS准备）
    extract_keywords: bool = True

    # PPTX参数
    slide_width_inches: float = 13.333
    slide_height_inches: float = 7.5
    
    def to_dict(self) -> Dict:
        return asdict(self)
    
    @classmethod
    def from_dict(cls, data: Dict) -> 'PipelineConfig':
        return cls(**{k: v for k, v in data.items() if k in cls.__dataclass_fields__})
    
    @classmethod
    def from_file(cls, path: str) -> 'PipelineConfig':
        with open(path, 'r', encoding='utf-8') as f:
            return cls.from_dict(json.load(f))
    
    def save(self, path: str):
        with open(path, 'w', encoding='utf-8') as f:
            json.dump(self.to_dict(), f, indent=2)


@dataclass
class BatchResult:
    """批量处理结果"""
    total: int
    success: int
    failed: int
    results: List[Dict]
    errors: List[Dict]


class SlidePipeline:
    """
    幻灯片处理流水线
    支持单页和批量处理，使用v4核心逻辑，DocLayout-YOLO版面检测 + VLM语义增强
    """

    def __init__(self, config: Optional[PipelineConfig] = None, use_vlm: bool = True, hybrid_mode: bool = True):
        self.config = config or PipelineConfig()
        self.use_vlm = use_vlm
        self.hybrid_mode = hybrid_mode if use_vlm else False  # 没有VLM时混合模式自动禁用

    def process_single(self, image_path: str, output_dir: str) -> Dict:
        """
        处理单张幻灯片（使用v4核心逻辑）

        Returns:
            包含处理结果的字典
        """
        try:
            # 调用新版v4流水线统一入口
            json_path, pptx_path = process_slide(
                image_path,
                output_dir,
                use_vlm=self.use_vlm,
                min_area=self.config.min_area,
                verbose=self.config.verbose,
                hybrid_mode=self.hybrid_mode
            )
            
            # 读取元数据
            with open(json_path, 'r', encoding='utf-8') as f:
                metadata = json.load(f)
            
            return {
                "status": "success",
                "image_path": image_path,
                "output_dir": os.path.dirname(json_path),
                "json_path": json_path,
                "pptx_path": pptx_path,
                "slide_id": metadata["slide_id"],
                "element_count": metadata["element_count"],
                "title": metadata.get("title", "")
            }
            
        except Exception as e:
            logger.error(f"处理失败 {image_path}: {e}")
            return {
                "status": "error",
                "image_path": image_path,
                "error": str(e)
            }
    
    def process_batch(self, image_paths: List[str], output_dir: str,
                      progress_callback: Optional[Callable] = None) -> BatchResult:
        """
        批量处理多张幻灯片
        
        Args:
            image_paths: 图像路径列表
            output_dir: 输出目录
            progress_callback: 进度回调函数 (current, total, result)
        """
        results = []
        errors = []
        
        with ThreadPoolExecutor(max_workers=self.config.max_workers) as executor:
            futures = {
                executor.submit(self.process_single, path, output_dir): path
                for path in image_paths
            }
            
            for i, future in enumerate(as_completed(futures)):
                result = future.result()
                
                if result["status"] == "success":
                    results.append(result)
                else:
                    errors.append(result)
                
                if progress_callback:
                    progress_callback(i + 1, len(image_paths), result)
        
        return BatchResult(
            total=len(image_paths),
            success=len(results),
            failed=len(errors),
            results=results,
            errors=errors
        )
    
    def process_directory(self, input_dir: str, output_dir: str,
                          patterns: List[str] = None) -> BatchResult:
        """
        处理目录中的所有图像
        
        Args:
            input_dir: 输入目录
            output_dir: 输出目录
            patterns: 文件模式列表，默认为 ['*.png', '*.jpg', '*.jpeg']
        """
        if patterns is None:
            patterns = ['*.png', '*.jpg', '*.jpeg', '*.PNG', '*.JPG', '*.JPEG']
        
        image_paths = []
        for pattern in patterns:
            image_paths.extend(glob.glob(os.path.join(input_dir, pattern)))
        
        image_paths = sorted(set(image_paths))
        
        if not image_paths:
            logger.warning(f"在 {input_dir} 中未找到图像文件")
            return BatchResult(0, 0, 0, [], [])
        
        logger.info(f"找到 {len(image_paths)} 个图像文件")
        
        def progress(current, total, result):
            status = "✓" if result["status"] == "success" else "✗"
            logger.info(f"[{current}/{total}] {status} {os.path.basename(result['image_path'])}")
        
        return self.process_batch(image_paths, output_dir, progress)


class PresentationBuilder:
    """
    演示文稿构建器
    从多个页面JSON创建多页PPTX
    """
    
    def __init__(self, slide_width_inches: float = 13.333, 
                 slide_height_inches: float = 7.5):
        self.prs = Presentation()
        self.prs.slide_width = Inches(slide_width_inches)
        self.prs.slide_height = Inches(slide_height_inches)
        self.slide_width = Inches(slide_width_inches)
        self.slide_height = Inches(slide_height_inches)
    
    def add_slide_from_json(self, json_path: str) -> int:
        """
        从JSON元数据添加幻灯片
        
        Returns:
            幻灯片索引
        """
        with open(json_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        
        metadata_dir = os.path.dirname(json_path)
        
        # 添加空白幻灯片
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        # 设置背景色
        bg_color = data.get('background_color', '#ffffff').lstrip('#')
        r, g, b = int(bg_color[0:2], 16), int(bg_color[2:4], 16), int(bg_color[4:6], 16)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(r, g, b)
        
        # 计算缩放比例
        original_width = data['width']
        original_height = data['height']
        scale_x = self.slide_width / original_width
        scale_y = self.slide_height / original_height
        
        # 添加元素（按z_order排序，保证层级正确）
        for elem_data in sorted(data['elements'], key=lambda e: e.get('z_order', 0)):
            bbox = elem_data['bbox']
            left = Emu(int(bbox['x'] * scale_x))
            top = Emu(int(bbox['y'] * scale_y))
            width = Emu(int(bbox['width'] * scale_x))
            height = Emu(int(bbox['height'] * scale_y))
            
            image_path = os.path.join(metadata_dir, elem_data['image_path'])
            if os.path.exists(image_path):
                shape = slide.shapes.add_picture(image_path, left, top, width, height)
                shape.name = elem_data['name']
        
        return len(self.prs.slides) - 1
    
    def add_slides_from_directory(self, directory: str) -> List[int]:
        """
        从目录中的所有JSON文件添加幻灯片
        按文件名排序
        """
        json_files = sorted(glob.glob(os.path.join(directory, "**/*.json"), recursive=True))
        indices = []
        
        for json_path in json_files:
            try:
                idx = self.add_slide_from_json(json_path)
                indices.append(idx)
                logger.info(f"添加幻灯片: {os.path.basename(json_path)}")
            except Exception as e:
                logger.error(f"添加幻灯片失败 {json_path}: {e}")
        
        return indices
    
    def save(self, output_path: str):
        """保存演示文稿"""
        self.prs.save(output_path)
        logger.info(f"演示文稿已保存: {output_path}")





# ============ 命令行接口 ============

def main():
    parser = argparse.ArgumentParser(
        description='幻灯片截图分割与PPTX还原流水线',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog='''
完整工作流程示例:
  1. 处理单张幻灯片
     python batch.py process slide.png -o ./output
  
  2. 批量处理目录
     python batch.py batch ./slides -o ./output
  
  3. 合并多个处理结果为单个PPTX
     python batch.py merge ./output -o presentation.pptx
  
  4. 生成解说旁白
     python batch.py narrate ./output -o narration.json
  
  5. 生成TTS语音和同步动画方案
     python batch.py tts narration.json -o ./tts_output
'''
    )
    
    subparsers = parser.add_subparsers(dest='command', help='可用命令')
    
    # 单张处理
    process_parser = subparsers.add_parser('process', help='处理单张幻灯片')
    process_parser.add_argument('image', help='输入图像路径')
    process_parser.add_argument('-o', '--output', default='./output', help='输出目录')
    process_parser.add_argument('--config', help='配置文件路径')
    process_parser.add_argument('--no-vlm', action='store_true', help='不使用VLM大模型（仅用传统CV/OCR fallback）')
    process_parser.add_argument('--no-hybrid', action='store_true', help='不使用混合模式（纯VLM直接检测边界，可能不稳定）')
    process_parser.add_argument('--fine-blocks', action='store_true', help='精细粒度模式（输出更多、更小的语义块，min_area=200）')
    process_parser.add_argument('--coarse-blocks', action='store_true', help='粗略粒度模式（输出更少、更大的语义块，min_area=400，推荐默认）')
    process_parser.add_argument('--fine', action='store_true', help='兼容旧参数，同--fine-blocks')
    process_parser.add_argument('--coarse', action='store_true', help='兼容旧参数，同--coarse-blocks')
    process_parser.add_argument('--min-area', type=int, help='自定义最小语义块面积（像素，默认300）')
    process_parser.add_argument('-v', '--verbose', action='store_true', help='输出详细调试信息')
    process_parser.add_argument('-q', '--quiet', action='store_true', help='静默模式，仅输出关键信息')

    # 批量处理
    batch_parser = subparsers.add_parser('batch', help='批量处理目录')
    batch_parser.add_argument('input_dir', help='输入目录')
    batch_parser.add_argument('-o', '--output', default='./output', help='输出目录')
    batch_parser.add_argument('--config', help='配置文件路径')
    batch_parser.add_argument('--workers', type=int, default=4, help='并行工作数')
    batch_parser.add_argument('--no-vlm', action='store_true', help='不使用VLM大模型（仅用传统CV/OCR fallback）')
    batch_parser.add_argument('--no-hybrid', action='store_true', help='不使用混合模式（纯VLM直接检测边界，可能不稳定）')
    batch_parser.add_argument('--fine-blocks', action='store_true', help='精细粒度模式（输出更多、更小的语义块，min_area=200）')
    batch_parser.add_argument('--coarse-blocks', action='store_true', help='粗略粒度模式（输出更少、更大的语义块，min_area=400，推荐默认）')
    batch_parser.add_argument('--fine', action='store_true', help='兼容旧参数，同--fine-blocks')
    batch_parser.add_argument('--coarse', action='store_true', help='兼容旧参数，同--coarse-blocks')
    batch_parser.add_argument('--min-area', type=int, help='自定义最小语义块面积（像素，默认300）')
    batch_parser.add_argument('-v', '--verbose', action='store_true', help='输出详细调试信息')
    batch_parser.add_argument('-q', '--quiet', action='store_true', help='静默模式，仅输出关键信息')
    
    # 合并PPTX
    merge_parser = subparsers.add_parser('merge', help='合并多个页面为单个PPTX')
    merge_parser.add_argument('input_dir', help='包含JSON文件的目录')
    merge_parser.add_argument('-o', '--output', default='merged.pptx', help='输出PPTX路径')
    
    # 生成旁白
    narrate_parser = subparsers.add_parser('narrate', help='生成旁白（为TTS准备）')
    narrate_parser.add_argument('input_dir', help='包含JSON文件的目录')
    narrate_parser.add_argument('-o', '--output', default='narration.json', help='输出JSON路径')
    
    tts_parser = subparsers.add_parser('tts', help='生成TTS语音和对应动画方案')
    tts_parser.add_argument('narration_json', help='旁白JSON文件路径')
    tts_parser.add_argument('-o', '--output', default='./tts_output', help='输出目录')
    tts_parser.add_argument('--voice', default='Cherry', help='发音人：Cherry(甜美女声)、Alvin(成熟男声)、Wanwan(可爱童声)')
    tts_parser.add_argument('--no-llm', action='store_true', help='不使用LLM，使用规则生成动画方案')
    
    args = parser.parse_args()
    
    if args.command == 'process':
        config = PipelineConfig.from_file(args.config) if args.config else PipelineConfig()
        # 命令行参数覆盖配置
        if args.fine or args.fine_blocks:
            config.min_area = args.min_area or 200
        elif args.coarse or args.coarse_blocks:
            config.min_area = args.min_area or 400
        elif args.min_area is not None:
            config.min_area = args.min_area

        # 输出级别设置
        if args.quiet:
            config.verbose = False
        elif args.verbose:
            config.verbose = True

        # 混合模式设置
        hybrid_mode = not args.no_hybrid if not args.no_vlm else False
        config.hybrid_mode = hybrid_mode

        pipeline = SlidePipeline(config, use_vlm=not args.no_vlm, hybrid_mode=hybrid_mode)
        result = pipeline.process_single(args.image, args.output)
        
        if result['status'] == 'success':
            print(f"✓ 处理成功")
            print(f"  JSON: {result['json_path']}")
            print(f"  PPTX: {result['pptx_path']}")
        else:
            print(f"✗ 处理失败: {result['error']}")
    
    elif args.command == 'batch':
        config = PipelineConfig.from_file(args.config) if args.config else PipelineConfig()
        config.max_workers = args.workers
        # 命令行参数覆盖配置
        if args.fine or args.fine_blocks:
            config.min_area = args.min_area or 200
        elif args.coarse or args.coarse_blocks:
            config.min_area = args.min_area or 400
        elif args.min_area is not None:
            config.min_area = args.min_area

        # 输出级别设置
        if args.quiet:
            config.verbose = False
        elif args.verbose:
            config.verbose = True

        # 混合模式设置
        hybrid_mode = not args.no_hybrid if not args.no_vlm else False
        config.hybrid_mode = hybrid_mode

        pipeline = SlidePipeline(config, use_vlm=not args.no_vlm, hybrid_mode=hybrid_mode)
        result = pipeline.process_directory(args.input_dir, args.output)
        
        print(f"\n处理完成:")
        print(f"  总计: {result.total}")
        print(f"  成功: {result.success}")
        print(f"  失败: {result.failed}")
    
    elif args.command == 'merge':
        builder = PresentationBuilder()
        builder.add_slides_from_directory(args.input_dir)
        builder.save(args.output)
        print(f"✓ 合并完成: {args.output}")
    
    elif args.command == 'narrate':
        narrations = generate_batch_narrations(args.input_dir, args.output)
        print(f"✓ 旁白生成完成: {args.output}")
        print(f"  共 {len(narrations)} 页")
    
    elif args.command == 'tts':
        result = generate_tts_and_animations(
            narration_json_path=args.narration_json,
            output_dir=args.output,
            voice=args.voice,
            use_llm_animation=not args.no_llm
        )
        print(f"✓ TTS生成和动画方案生成完成")
        print(f"  音频输出目录: {args.output}")
        print(f"  使用发音人: {args.voice}")
        print(f"  动画生成方式: {'LLM智能生成' if not args.no_llm else '规则生成'}")

        # 兼容批量和单页结果
        if result.get("type") == "batch":
            total_duration = sum([r["audio_info"]["total_duration"] for r in result["results"]])
            total_animations = sum([len(r["animation_scheme"]["animations"]) for r in result["results"]])
            print(f"  批量处理幻灯片数量: {result['count']} 个")
            print(f"  所有旁白总时长: {total_duration:.2f} 秒")
            print(f"  总动画数量: {total_animations} 个")
        else:
            print(f"  总旁白时长: {result['audio_info']['total_duration']:.2f} 秒")
            print(f"  总动画时长: {result['animation_scheme']['total_duration']:.2f} 秒")
            print(f"  动画数量: {len(result['animation_scheme']['animations'])} 个")
    
    else:
        parser.print_help()


if __name__ == "__main__":
    main()
