"""
幻灯片PPTX重建器
Slide Reconstructor - PPTX Reconstruction

从元数据JSON重建PPTX文件，支持背景遮罩、元素定位插入
"""

import os
import json
from typing import List, Dict

import cv2
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.dml.color import RGBColor


class SlideReconstructor:
    """幻灯片重建器，生成PPTX文件"""

    def __init__(self,
                 slide_width_inches: float = 13.333,
                 slide_height_inches: float = 7.5,
                 use_original_as_background: bool = True,
                 mask_elements: bool = True):
        self.slide_width = Inches(slide_width_inches)
        self.slide_height = Inches(slide_height_inches)
        self.use_original_as_background = use_original_as_background
        self.mask_elements = mask_elements

    def reconstruct(self, metadata_path: str, output_path: str) -> str:
        """从元数据重建PPTX"""
        with open(metadata_path, 'r', encoding='utf-8') as f:
            data = json.load(f)

        metadata_dir = os.path.dirname(metadata_path)

        # 创建PPTX
        prs = Presentation()
        prs.slide_width = self.slide_width
        prs.slide_height = self.slide_height
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # 计算缩放比例
        original_width = data['width']
        original_height = data['height']
        scale_x = self.slide_width / original_width
        scale_y = self.slide_height / original_height

        bg_color = data.get('background_color', '#ffffff')

        # 设置背景
        if self.use_original_as_background:
            original_image_path = None
            for fname in os.listdir(metadata_dir):
                if fname.startswith("original_"):
                    original_image_path = os.path.join(metadata_dir, fname)
                    break

            if original_image_path and os.path.exists(original_image_path):
                if self.mask_elements:
                    masked_bg_path = self._create_masked_background(
                        original_image_path, data['elements'], bg_color, metadata_dir
                    )
                    slide.shapes.add_picture(masked_bg_path, 0, 0, self.slide_width, self.slide_height)
                else:
                    slide.shapes.add_picture(original_image_path, 0, 0, self.slide_width, self.slide_height)
            else:
                self._set_background_color(slide, bg_color)
        else:
            self._set_background_color(slide, bg_color)

        # 插入元素（按z_order排序）
        elements = sorted(data['elements'], key=lambda e: e.get('z_order', 0))
        for elem_data in elements:
            self._add_element(slide, elem_data, metadata_dir, scale_x, scale_y)

        prs.save(output_path)
        return output_path

    def _create_masked_background(self, original_path: str, elements: List[Dict],
                                  bg_color: str, output_dir: str) -> str:
        """创建元素位置遮罩后的背景图"""
        image = cv2.imread(original_path)
        image_rgb = cv2.cvtColor(image, cv2.COLOR_BGR2RGB)

        bg_hex = bg_color.lstrip('#')
        bg_rgb = (int(bg_hex[0:2], 16), int(bg_hex[2:4], 16), int(bg_hex[4:6], 16))

        for elem in elements:
            bbox = elem['bbox']
            x, y = bbox['x'], bbox['y']
            w, h = bbox['width'], bbox['height']

            padding = 2
            x1 = max(0, x - padding)
            y1 = max(0, y - padding)
            x2 = min(image_rgb.shape[1], x + w + padding)
            y2 = min(image_rgb.shape[0], y + h + padding)

            image_rgb[y1:y2, x1:x2] = bg_rgb

        masked_path = os.path.join(output_dir, "masked_background.png")
        Image.fromarray(image_rgb).save(masked_path)
        return masked_path

    def _set_background_color(self, slide, hex_color: str):
        """设置纯色背景"""
        hex_color = hex_color.lstrip('#')
        r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(r, g, b)

    def _add_element(self, slide, elem_data: Dict, base_dir: str, scale_x: Emu, scale_y: Emu):
        """添加元素到PPT"""
        bbox = elem_data['bbox']
        left = int(bbox['x'] * scale_x)
        top = int(bbox['y'] * scale_y)
        width = int(bbox['width'] * scale_x)
        height = int(bbox['height'] * scale_y)

        image_path = os.path.join(base_dir, elem_data['image_path'])
        if os.path.exists(image_path):
            shape = slide.shapes.add_picture(image_path, left, top, width, height)
            shape.name = elem_data['name']
        else:
            print(f"警告: 元素图片不存在: {image_path}")
