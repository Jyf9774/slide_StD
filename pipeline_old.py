#!/usr/bin/env python3
"""
幻灯片截图分割与PPTX还原数据流水线 v3.0
Slide Screenshot Segmentation and PPTX Reconstruction Pipeline

核心优化：
1. 使用GPT-4.1视觉能力替代传统OCR和规则分类
2. 原图作为背景，元素位置用背景色遮罩后放置提取元素
3. 简化CV流程，仅保留边界框检测
"""

import os
import json
import hashlib
import base64
from pathlib import Path
from dataclasses import dataclass, field
from typing import List, Tuple, Optional, Dict, Any
from enum import Enum
from datetime import datetime

import cv2
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ============ GPT-4.1 集成 ============

def get_gpt_client():
    """获取Azure OpenAI客户端"""
    try:
        # 尝试导入openai模块
        try:
            from openai import AzureOpenAI
            print("✅ 成功导入openai模块")
        except ImportError:
            print("❌ 未找到openai模块，请安装: pip install openai")
            return None
        
        # 使用更可靠的方式导入key_manage
        try:
            import sys
            import os
            
            # 确保当前目录在Python路径中
            if os.path.dirname(os.path.abspath(__file__)) not in sys.path:
                sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
                
            from key_manage import Azure_GPT_4_1_Key
            print("✅ 成功导入key_manage模块")
        except ImportError as ie:
            print(f"❌ 导入key_manage失败: {str(ie)}")
            print(f"   当前目录: {os.getcwd()}")
            print(f"   文件路径: {os.path.abspath(__file__)}")
            return None
        
        # 初始化客户端
        try:
            client = AzureOpenAI(
                api_version="2024-12-01-preview",
                api_key=Azure_GPT_4_1_Key,
                azure_endpoint="https://admin-m9uwcx36-eastus2.cognitiveservices.azure.com/"
            )
            print("✅ 成功创建Azure OpenAI客户端")
            return client
        except Exception as e:
            print(f"❌ Azure OpenAI客户端初始化失败: {str(e)}")
            return None
    except Exception as e:
        print(f"❌ GPT客户端初始化发生未知错误: {str(e)}")
        return None


def image_to_base64(image_path: str) -> str:
    """将图片转换为base64"""
    with open(image_path, "rb") as f:
        return base64.b64encode(f.read()).decode("utf-8")


def numpy_to_base64(image_array: np.ndarray) -> str:
    """将numpy数组转换为base64"""
    # RGB转BGR（OpenCV格式）
    if len(image_array.shape) == 3 and image_array.shape[2] == 3:
        image_bgr = cv2.cvtColor(image_array, cv2.COLOR_RGB2BGR)
    else:
        image_bgr = image_array
    
    _, buffer = cv2.imencode('.png', image_bgr)
    return base64.b64encode(buffer).decode("utf-8")


class GPTAnalyzer:
    """GPT-4.1 视觉分析器"""
    
    def __init__(self, client=None):
        self.client = client or get_gpt_client()
        self.enabled = self.client is not None
    
    def analyze_slide(self, image_path: str) -> Dict:
        """
        分析整张幻灯片
        返回：背景色、整体描述、标题等
        """
        if not self.enabled:
            return {}
        
        base64_image = image_to_base64(image_path)
        
        prompt = """请分析这张幻灯片截图，返回JSON格式的分析结果：

{
    "background_color": "#RRGGBB格式的背景色，优先识别主要背景区域的颜色",
    "title": "幻灯片的主标题文本",
    "subtitle": "副标题（如果有）",
    "slide_type": "封面页/内容页/图表页/结束页 等",
    "main_language": "中文/英文/中英混合",
    "description": "用1-2句话描述这张幻灯片的主要内容",
    "key_points": ["要点1", "要点2", "..."]
}

注意：
1. 背景色请识别幻灯片的主要背景区域（通常是白色或浅色），忽略边缘装饰
2. 只返回JSON，不要其他文字"""

        try:
            resp = self.client.chat.completions.create(
                model="gpt-4.1",
                messages=[
                    {"role": "system", "content": "你是一个专业的幻灯片分析助手，擅长识别和提取幻灯片中的结构化信息。"},
                    {"role": "user", "content": [
                        {"type": "text", "text": prompt},
                        {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{base64_image}"}}
                    ]}
                ],
                temperature=0.1,
                max_tokens=1000
            )
            
            result_text = resp.choices[0].message.content.strip()
            # 清理可能的markdown代码块
            if result_text.startswith("```"):
                result_text = result_text.split("```")[1]
                if result_text.startswith("json"):
                    result_text = result_text[4:]
            result_text = result_text.strip()
            
            return json.loads(result_text)
        except Exception as e:
            print(f"GPT分析幻灯片失败: {e}")
            return {}
    
    def analyze_element(self, element_image: np.ndarray, context: str = "") -> Dict:
        """
        分析单个元素区域
        返回：类型、文本内容、描述等
        """
        if not self.enabled:
            return {}
        
        base64_image = numpy_to_base64(element_image)
        
        prompt = f"""请分析这个幻灯片元素区域，返回JSON格式的分析结果：

{{
    "type": "text/image/chart/table/logo/diagram/decoration",
    "text_content": "提取所有可见文本，保持原始换行格式",
    "description": "简短描述这个元素的内容",
    "is_title": true/false,
    "estimated_font_size": 数字（估计的主要字体大小，单位pt）,
    "text_color": "#RRGGBB格式的主要文字颜色",
    "has_border": true/false
}}

{f'上下文信息：{context}' if context else ''}

类型说明：
- text: 纯文本区域
- image: 图片/照片
- chart: 图表（柱状图、折线图等）
- table: 表格
- logo: 标志/图标
- diagram: 示意图/流程图
- decoration: 装饰元素

只返回JSON，不要其他文字。"""

        try:
            resp = self.client.chat.completions.create(
                model="gpt-4.1",
                messages=[
                    {"role": "system", "content": "你是一个OCR和图像分析专家，擅长识别图像中的文字和元素类型。"},
                    {"role": "user", "content": [
                        {"type": "text", "text": prompt},
                        {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{base64_image}"}}
                    ]}
                ],
                temperature=0.1,
                max_tokens=800
            )
            
            result_text = resp.choices[0].message.content.strip()
            if result_text.startswith("```"):
                result_text = result_text.split("```")[1]
                if result_text.startswith("json"):
                    result_text = result_text[4:]
            result_text = result_text.strip()
            
            return json.loads(result_text)
        except Exception as e:
            print(f"GPT分析元素失败: {e}")
            return {}
    
    def batch_analyze_elements(self, elements_data: List[Tuple[np.ndarray, str]]) -> List[Dict]:
        """批量分析元素（减少API调用次数）"""
        results = []
        for img, context in elements_data:
            result = self.analyze_element(img, context)
            results.append(result)
        return results


# ============ 传统方法备用 ============

class FallbackAnalyzer:
    """传统CV/OCR分析器（GPT不可用时的备用方案）"""
    
    def __init__(self, ocr_lang: str = "chi_sim+eng"):
        self.ocr_lang = ocr_lang
        self._init_tesseract()
    
    def _init_tesseract(self):
        """初始化Tesseract"""
        try:
            import pytesseract
            tesseract_paths = ['/opt/homebrew/bin/tesseract', '/usr/local/bin/tesseract', '/usr/bin/tesseract']
            for path in tesseract_paths:
                if os.path.exists(path):
                    pytesseract.pytesseract.tesseract_cmd = path
                    break
            self.pytesseract = pytesseract
            self.ocr_available = True
        except ImportError:
            self.ocr_available = False
            print("警告: pytesseract未安装，OCR功能不可用")
    
    def extract_background_color(self, image_rgb: np.ndarray) -> str:
        """提取背景色（白色优先策略）"""
        h, w = image_rgb.shape[:2]
        
        # 采样中心区域
        margin_x, margin_y = int(w * 0.15), int(h * 0.15)
        center = image_rgb[margin_y:h-margin_y, margin_x:w-margin_x]
        
        if center.size == 0:
            center = image_rgb
        
        # 白色检测
        white_mask = np.all(center > 240, axis=2)
        white_ratio = np.sum(white_mask) / white_mask.size
        
        if white_ratio > 0.3:
            return "#ffffff"
        
        # 浅色检测
        light_mask = np.all(center > 200, axis=2)
        light_ratio = np.sum(light_mask) / light_mask.size
        
        if light_ratio > 0.4:
            light_pixels = center[light_mask]
            if len(light_pixels) > 0:
                avg = np.mean(light_pixels, axis=0).astype(int)
                return "#{:02x}{:02x}{:02x}".format(*avg)
        
        # K-means聚类
        try:
            from sklearn.cluster import KMeans
            pixels = center.reshape(-1, 3)
            sample = pixels[np.random.choice(len(pixels), min(5000, len(pixels)), replace=False)]
            kmeans = KMeans(n_clusters=3, random_state=42, n_init=10).fit(sample)
            labels, counts = np.unique(kmeans.labels_, return_counts=True)
            dominant = kmeans.cluster_centers_[labels[np.argmax(counts)]].astype(int)
            return "#{:02x}{:02x}{:02x}".format(*dominant)
        except:
            return "#ffffff"
    
    def ocr_region(self, region: np.ndarray) -> str:
        """OCR识别"""
        if not self.ocr_available:
            return ""
        
        try:
            enlarged = cv2.resize(region, None, fx=2, fy=2, interpolation=cv2.INTER_CUBIC)
            pil_image = Image.fromarray(enlarged)
            text = self.pytesseract.image_to_string(pil_image, lang=self.ocr_lang, config='--oem 3 --psm 6')
            return '\n'.join([l.strip() for l in text.split('\n') if l.strip()])
        except Exception as e:
            print(f"OCR错误: {e}")
            return ""
    
    def classify_region(self, region: np.ndarray) -> str:
        """分类区域类型"""
        h, w = region.shape[:2]
        gray = cv2.cvtColor(region, cv2.COLOR_RGB2GRAY)
        
        edges = cv2.Canny(gray, 50, 150)
        edge_density = np.sum(edges > 0) / edges.size
        
        hsv = cv2.cvtColor(region, cv2.COLOR_RGB2HSV)
        h_std, s_std = np.std(hsv[:,:,0]), np.std(hsv[:,:,1])
        
        lines = cv2.HoughLinesP(edges, 1, np.pi/180, 50, minLineLength=30, maxLineGap=10)
        line_count = len(lines) if lines is not None else 0
        
        white_ratio = np.sum(np.all(region > 220, axis=2)) / (h * w)
        
        if w / h > 10 or h / w > 10:
            return "decoration"
        if line_count > 15:
            return "table"
        if line_count > 5 and s_std > 40:
            return "chart"
        if white_ratio > 0.5 and edge_density < 0.15:
            return "text"
        if h_std > 40 and s_std > 50:
            return "image"
        return "mixed"


# ============ 数据结构 ============

class NumpyEncoder(json.JSONEncoder):
    def default(self, obj):
        if isinstance(obj, (np.integer, np.floating, np.bool_)):
            return obj.item()
        if isinstance(obj, np.ndarray):
            return obj.tolist()
        return super().default(obj)


class ElementType(Enum):
    TEXT = "text"
    IMAGE = "image"
    CHART = "chart"
    TABLE = "table"
    DIAGRAM = "diagram"
    LOGO = "logo"
    MIXED = "mixed"
    DECORATION = "decoration"
    UNKNOWN = "unknown"


@dataclass
class BoundingBox:
    x: int
    y: int
    width: int
    height: int
    
    @property
    def x2(self) -> int:
        return self.x + self.width
    
    @property
    def y2(self) -> int:
        return self.y + self.height
    
    @property
    def area(self) -> int:
        return self.width * self.height
    
    def to_dict(self) -> Dict:
        return {"x": self.x, "y": self.y, "width": self.width, "height": self.height}


@dataclass 
class SlideElement:
    id: str
    name: str
    type: ElementType
    bbox: BoundingBox
    image_path: str
    text_content: str = ""
    confidence: float = 0.0
    dominant_colors: List[str] = field(default_factory=list)
    has_border: bool = False
    is_title: bool = False
    z_order: int = 0
    font_size_estimate: int = 12
    text_color: str = "#000000"
    description: str = ""
    metadata: Dict[str, Any] = field(default_factory=dict)
    
    def to_dict(self) -> Dict:
        return {
            "id": self.id,
            "name": self.name,
            "type": self.type.value,
            "bbox": self.bbox.to_dict(),
            "image_path": self.image_path,
            "text_content": self.text_content,
            "confidence": self.confidence,
            "dominant_colors": self.dominant_colors,
            "has_border": self.has_border,
            "is_title": self.is_title,
            "z_order": self.z_order,
            "font_size_estimate": self.font_size_estimate,
            "text_color": self.text_color,
            "description": self.description,
            "metadata": self.metadata
        }


@dataclass
class SlideMetadata:
    slide_id: str
    source_image: str
    width: int
    height: int
    aspect_ratio: str
    background_color: str
    background_colors_palette: List[str]
    element_count: int
    elements: List[SlideElement]
    created_at: str
    version: str = "3.0"
    title: str = ""
    subtitle: str = ""
    slide_type: str = ""
    description: str = ""
    key_points: List[str] = field(default_factory=list)
    keywords: List[str] = field(default_factory=list)
    
    def to_dict(self) -> Dict:
        return {
            "slide_id": self.slide_id,
            "source_image": self.source_image,
            "width": self.width,
            "height": self.height,
            "aspect_ratio": self.aspect_ratio,
            "background_color": self.background_color,
            "background_colors_palette": self.background_colors_palette,
            "element_count": self.element_count,
            "elements": [e.to_dict() for e in self.elements],
            "created_at": self.created_at,
            "version": self.version,
            "title": self.title,
            "subtitle": self.subtitle,
            "slide_type": self.slide_type,
            "description": self.description,
            "key_points": self.key_points,
            "keywords": self.keywords
        }


# ============ 核心流水线 ============

class SlideSegmenter:
    """幻灯片分割器"""
    
    def __init__(self, 
                 min_area: int = 500,
                 merge_threshold: int = 15,
                 use_gpt: bool = True):
        self.min_area = min_area
        self.merge_threshold = merge_threshold
        self.use_gpt = use_gpt
        
        # 初始化分析器
        self.gpt_analyzer = GPTAnalyzer() if use_gpt else None
        self.fallback_analyzer = FallbackAnalyzer()
        
        # 判断是否可以使用GPT
        self.gpt_available = self.gpt_analyzer and self.gpt_analyzer.enabled
        if use_gpt and not self.gpt_available:
            print("提示: GPT-4.1不可用，将使用传统CV/OCR方法")
        
    def segment(self, image_path: str, output_dir: str) -> SlideMetadata:
        """分割幻灯片图像"""
        output_path = Path(output_dir)
        output_path.mkdir(parents=True, exist_ok=True)
        elements_dir = output_path / "elements"
        elements_dir.mkdir(exist_ok=True)
        
        # 读取图像
        image = cv2.imread(image_path)
        if image is None:
            raise ValueError(f"无法读取图像: {image_path}")
        
        height, width = image.shape[:2]
        image_rgb = cv2.cvtColor(image, cv2.COLOR_BGR2RGB)
        
        slide_id = self._generate_id(image_path)
        
        # ====== 使用GPT分析整体幻灯片 ======
        slide_analysis = {}
        if self.gpt_available:
            print("正在使用GPT-4.1分析幻灯片...")
            slide_analysis = self.gpt_analyzer.analyze_slide(image_path)
        
        # 提取背景色
        if slide_analysis.get('background_color'):
            bg_color = slide_analysis['background_color']
        else:
            bg_color = self.fallback_analyzer.extract_background_color(image_rgb)
        
        bg_palette = [bg_color]
        
        # ====== CV检测区域边界框 ======
        regions = self._detect_regions(image)
        merged_regions = self._merge_overlapping_regions(regions)
        
        # ====== 处理每个区域 ======
        elements = []
        for idx, bbox in enumerate(merged_regions):
            element = self._process_region(
                image_rgb, bbox, idx, 
                str(elements_dir), slide_id,
                slide_analysis.get('title', '')
            )
            if element:
                elements.append(element)
        
        # 排序
        elements.sort(key=lambda e: (e.bbox.y, e.bbox.x))
        for idx, elem in enumerate(elements):
            elem.z_order = idx
        
        # 创建元数据
        metadata = SlideMetadata(
            slide_id=slide_id,
            source_image=os.path.basename(image_path),
            width=width,
            height=height,
            aspect_ratio=self._calc_aspect_ratio(width, height),
            background_color=bg_color,
            background_colors_palette=bg_palette,
            element_count=len(elements),
            elements=elements,
            created_at=datetime.now().isoformat(),
            title=slide_analysis.get('title', ''),
            subtitle=slide_analysis.get('subtitle', ''),
            slide_type=slide_analysis.get('slide_type', ''),
            description=slide_analysis.get('description', ''),
            key_points=slide_analysis.get('key_points', [])
        )
        
        # 提取关键词
        all_text = " ".join([e.text_content for e in elements if e.text_content])
        metadata.keywords = list(set([w for w in all_text.split() if len(w) > 2]))[:30]
        
        # 保存JSON
        json_path = output_path / f"{slide_id}.json"
        with open(json_path, 'w', encoding='utf-8') as f:
            json.dump(metadata.to_dict(), f, ensure_ascii=False, indent=2, cls=NumpyEncoder)
        
        # 复制原图
        original_copy = output_path / f"original_{os.path.basename(image_path)}"
        cv2.imwrite(str(original_copy), image)
        
        return metadata
    
    def _generate_id(self, image_path: str) -> str:
        content = f"{image_path}_{datetime.now().isoformat()}"
        return hashlib.md5(content.encode()).hexdigest()[:12]
    
    def _detect_regions(self, image: np.ndarray) -> List[BoundingBox]:
        """CV检测区域边界框"""
        gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
        h, w = image.shape[:2]
        
        # 自适应阈值 + Canny边缘
        binary = cv2.adaptiveThreshold(gray, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
                                        cv2.THRESH_BINARY_INV, 11, 2)
        edges = cv2.Canny(gray, 50, 150)
        combined = cv2.bitwise_or(binary, edges)
        
        # 形态学操作
        kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (20, 8))
        dilated = cv2.dilate(combined, kernel, iterations=2)
        closed = cv2.morphologyEx(dilated, cv2.MORPH_CLOSE, 
                                   cv2.getStructuringElement(cv2.MORPH_RECT, (10, 10)))
        
        contours, _ = cv2.findContours(closed, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
        
        regions = []
        for contour in contours:
            x, y, cw, ch = cv2.boundingRect(contour)
            area = cw * ch
            
            if area < self.min_area or area > 0.95 * w * h:
                continue
            
            aspect = cw / ch if ch > 0 else 0
            if aspect > 20 or aspect < 0.05:
                continue
            
            padding = 8
            x = max(0, x - padding)
            y = max(0, y - padding)
            cw = min(w - x, cw + 2 * padding)
            ch = min(h - y, ch + 2 * padding)
            
            regions.append(BoundingBox(x, y, cw, ch))
        
        return regions
    
    def _merge_overlapping_regions(self, regions: List[BoundingBox]) -> List[BoundingBox]:
        """合并重叠区域"""
        if not regions:
            return []
        
        boxes = [[r.x, r.y, r.x2, r.y2] for r in regions]
        
        merged = True
        while merged:
            merged = False
            new_boxes = []
            used = set()
            
            for i, box1 in enumerate(boxes):
                if i in used:
                    continue
                x1, y1, x2, y2 = box1
                
                for j, box2 in enumerate(boxes[i+1:], i+1):
                    if j in used:
                        continue
                    bx1, by1, bx2, by2 = box2
                    
                    if (x1 - self.merge_threshold <= bx2 and 
                        x2 + self.merge_threshold >= bx1 and
                        y1 - self.merge_threshold <= by2 and 
                        y2 + self.merge_threshold >= by1):
                        x1, y1 = min(x1, bx1), min(y1, by1)
                        x2, y2 = max(x2, bx2), max(y2, by2)
                        used.add(j)
                        merged = True
                
                new_boxes.append([x1, y1, x2, y2])
                used.add(i)
            
            boxes = new_boxes
        
        return [BoundingBox(b[0], b[1], b[2]-b[0], b[3]-b[1]) for b in boxes]
    
    def _process_region(self, image_rgb: np.ndarray, bbox: BoundingBox, 
                        idx: int, output_dir: str, slide_id: str,
                        slide_title: str = "") -> Optional[SlideElement]:
        """处理单个区域"""
        crop = image_rgb[bbox.y:bbox.y2, bbox.x:bbox.x2]
        
        if crop.size == 0:
            return None
        
        element_id = f"{slide_id}_elem_{idx:03d}"
        element_name = f"element_{idx:03d}"
        
        # 保存裁剪图像
        image_filename = f"{element_name}.png"
        image_path = os.path.join(output_dir, image_filename)
        Image.fromarray(crop).save(image_path)
        
        # ====== 使用GPT或传统方法分析元素 ======
        if self.gpt_available:
            context = f"这是幻灯片'{slide_title}'中的一个区域" if slide_title else ""
            analysis = self.gpt_analyzer.analyze_element(crop, context)
            
            elem_type_str = analysis.get('type', 'mixed')
            text_content = analysis.get('text_content', '')
            description = analysis.get('description', '')
            is_title = analysis.get('is_title', False)
            font_size = analysis.get('estimated_font_size', 12)
            text_color = analysis.get('text_color', '#000000')
            has_border = analysis.get('has_border', False)
            confidence = 0.9
        else:
            elem_type_str = self.fallback_analyzer.classify_region(crop)
            text_content = self.fallback_analyzer.ocr_region(crop) if elem_type_str in ['text', 'mixed', 'table'] else ""
            description = ""
            is_title = bbox.y < image_rgb.shape[0] * 0.25 and bbox.width > image_rgb.shape[1] * 0.3
            font_size = int(bbox.height / max(1, len(text_content.split('\n'))) * 0.6) if text_content else 12
            font_size = max(8, min(72, font_size))
            text_color = "#000000"
            has_border = self._detect_border(crop)
            confidence = 0.6
        
        # 转换类型枚举
        try:
            elem_type = ElementType(elem_type_str.lower())
        except ValueError:
            elem_type = ElementType.MIXED
        
        # 提取主要颜色
        dominant_colors = self._get_dominant_colors(crop)
        
        return SlideElement(
            id=element_id,
            name=element_name,
            type=elem_type,
            bbox=bbox,
            image_path=f"elements/{image_filename}",
            text_content=text_content,
            confidence=confidence,
            dominant_colors=dominant_colors,
            has_border=has_border,
            is_title=is_title,
            font_size_estimate=font_size,
            text_color=text_color,
            description=description,
            metadata={
                "aspect_ratio": round(bbox.width / bbox.height, 2) if bbox.height > 0 else 0,
                "relative_size": bbox.area
            }
        )
    
    def _get_dominant_colors(self, region: np.ndarray, n_colors: int = 3) -> List[str]:
        """提取主要颜色"""
        try:
            from sklearn.cluster import KMeans
            pixels = region.reshape(-1, 3)
            sample = pixels[np.random.choice(len(pixels), min(3000, len(pixels)), replace=False)]
            kmeans = KMeans(n_clusters=min(n_colors, len(sample)//10+1), random_state=42, n_init=10).fit(sample)
            return ["#{:02x}{:02x}{:02x}".format(*c.astype(int)) for c in kmeans.cluster_centers_]
        except:
            return []
    
    def _detect_border(self, region: np.ndarray) -> bool:
        """检测边框"""
        gray = cv2.cvtColor(region, cv2.COLOR_RGB2GRAY)
        edges = cv2.Canny(gray, 50, 150)
        h, w = edges.shape
        bw = 3
        
        scores = [
            np.sum(edges[:bw, :] > 0) / (bw * w) if w > 0 else 0,
            np.sum(edges[-bw:, :] > 0) / (bw * w) if w > 0 else 0,
            np.sum(edges[:, :bw] > 0) / (h * bw) if h > 0 else 0,
            np.sum(edges[:, -bw:] > 0) / (h * bw) if h > 0 else 0
        ]
        return sum(s > 0.3 for s in scores) >= 3
    
    def _calc_aspect_ratio(self, width: int, height: int) -> str:
        from math import gcd
        d = gcd(width, height)
        return f"{width//d}:{height//d}"


class SlideReconstructor:
    """幻灯片重建器 - 原图背景 + 元素遮罩方案"""
    
    def __init__(self, 
                 slide_width_inches: float = 13.333, 
                 slide_height_inches: float = 7.5,
                 use_original_as_background: bool = True,
                 mask_elements: bool = True):
        """
        Args:
            use_original_as_background: 使用原图作为背景
            mask_elements: 在元素位置用背景色遮罩
        """
        self.slide_width = Inches(slide_width_inches)
        self.slide_height = Inches(slide_height_inches)
        self.use_original_as_background = use_original_as_background
        self.mask_elements = mask_elements
    
    def reconstruct(self, metadata_path: str, output_path: str) -> str:
        """从元数据重建幻灯片"""
        with open(metadata_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        
        metadata_dir = os.path.dirname(metadata_path)
        
        prs = Presentation()
        prs.slide_width = self.slide_width
        prs.slide_height = self.slide_height
        
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        # 计算缩放比例
        original_width = data['width']
        original_height = data['height']
        scale_x = self.slide_width / original_width
        scale_y = self.slide_height / original_height
        
        bg_color = data.get('background_color', '#ffffff')
        
        # ====== 方案：原图背景 + 遮罩 ======
        if self.use_original_as_background:
            # 1. 查找原图
            original_image_path = None
            for fname in os.listdir(metadata_dir):
                if fname.startswith("original_"):
                    original_image_path = os.path.join(metadata_dir, fname)
                    break
            
            if original_image_path and os.path.exists(original_image_path):
                if self.mask_elements:
                    # 创建带遮罩的背景图
                    masked_bg_path = self._create_masked_background(
                        original_image_path, data['elements'], bg_color, metadata_dir
                    )
                    # 插入遮罩后的背景图
                    slide.shapes.add_picture(masked_bg_path, 0, 0, self.slide_width, self.slide_height)
                else:
                    # 直接插入原图作为背景
                    slide.shapes.add_picture(original_image_path, 0, 0, self.slide_width, self.slide_height)
            else:
                # 没有原图，使用纯色背景
                self._set_background_color(slide, bg_color)
        else:
            self._set_background_color(slide, bg_color)
        
        # ====== 插入元素 ======
        elements = sorted(data['elements'], key=lambda e: e.get('z_order', 0))
        
        for elem_data in elements:
            self._add_element(slide, elem_data, metadata_dir, scale_x, scale_y)
        
        prs.save(output_path)
        return output_path
    
    def _create_masked_background(self, original_path: str, elements: List[Dict], 
                                   bg_color: str, output_dir: str) -> str:
        """创建带遮罩的背景图"""
        # 读取原图
        image = cv2.imread(original_path)
        image_rgb = cv2.cvtColor(image, cv2.COLOR_BGR2RGB)
        
        # 解析背景色
        bg_hex = bg_color.lstrip('#')
        bg_rgb = (int(bg_hex[0:2], 16), int(bg_hex[2:4], 16), int(bg_hex[4:6], 16))
        
        # 在元素位置绘制背景色矩形
        for elem in elements:
            bbox = elem['bbox']
            x, y = bbox['x'], bbox['y']
            w, h = bbox['width'], bbox['height']
            
            # 稍微扩大遮罩区域，确保完全覆盖
            padding = 2
            x1 = max(0, x - padding)
            y1 = max(0, y - padding)
            x2 = min(image_rgb.shape[1], x + w + padding)
            y2 = min(image_rgb.shape[0], y + h + padding)
            
            # 填充背景色
            image_rgb[y1:y2, x1:x2] = bg_rgb
        
        # 保存遮罩后的图像
        masked_path = os.path.join(output_dir, "masked_background.png")
        Image.fromarray(image_rgb).save(masked_path)
        
        return masked_path
    
    def _set_background_color(self, slide, hex_color: str):
        """设置幻灯片背景色"""
        hex_color = hex_color.lstrip('#')
        r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(r, g, b)
    
    def _add_element(self, slide, elem_data: Dict, base_dir: str, 
                     scale_x: Emu, scale_y: Emu):
        """添加元素"""
        bbox = elem_data['bbox']
        
        left = int(bbox['x'] * scale_x)
        top = int(bbox['y'] * scale_y)
        width = int(bbox['width'] * scale_x)
        height = int(bbox['height'] * scale_y)
        
        image_path = os.path.join(base_dir, elem_data['image_path'])
        
        if os.path.exists(image_path):
            shape = slide.shapes.add_picture(image_path, left, top, width, height)
            shape.name = elem_data['name']
        else:
            print(f"警告: 图片不存在 {image_path}")


# ============ 便捷函数 ============

def process_slide(image_path: str, 
                  output_base_dir: str = "./output",
                  use_gpt: bool = True,
                  use_original_bg: bool = True,
                  mask_elements: bool = True) -> Tuple[str, str]:
    """
    处理单张幻灯片
    
    Args:
        image_path: 输入图像路径
        output_base_dir: 输出基础目录
        use_gpt: 是否使用GPT-4.1进行分析
        use_original_bg: 是否使用原图作为背景
        mask_elements: 是否在元素位置遮罩
        
    Returns:
        (json_path, pptx_path): JSON和PPTX文件路径
    """
    image_name = os.path.splitext(os.path.basename(image_path))[0]
    output_dir = os.path.join(output_base_dir, image_name)
    
    # 分割
    segmenter = SlideSegmenter(use_gpt=use_gpt)
    metadata = segmenter.segment(image_path, output_dir)
    
    json_path = os.path.join(output_dir, f"{metadata.slide_id}.json")
    
    # 重建
    pptx_path = os.path.join(output_dir, f"{image_name}_reconstructed.pptx")
    reconstructor = SlideReconstructor(
        use_original_as_background=use_original_bg,
        mask_elements=mask_elements
    )
    reconstructor.reconstruct(json_path, pptx_path)
    
    return json_path, pptx_path


if __name__ == "__main__":
    import sys
    
    print("""
幻灯片分割与重建流水线 v3.0
===========================
用法: python pipeline.py <image_path> [output_dir] [options]

选项:
  --no-gpt        不使用GPT-4.1（使用传统CV/OCR）
  --no-bg         不使用原图作为背景
  --no-mask       不遮罩元素区域

示例:
  python pipeline.py slide.png ./output
  python pipeline.py slide.png ./output --no-gpt
""")
    
    if len(sys.argv) < 2:
        sys.exit(1)
    
    image_path = sys.argv[1]
    output_dir = "./output"
    use_gpt = True
    use_original_bg = True
    mask_elements = True
    
    for arg in sys.argv[2:]:
        if arg == "--no-gpt":
            use_gpt = False
        elif arg == "--no-bg":
            use_original_bg = False
        elif arg == "--no-mask":
            mask_elements = False
        elif not arg.startswith("-"):
            output_dir = arg
    
    print(f"处理图像: {image_path}")
    print(f"输出目录: {output_dir}")
    print(f"使用GPT: {use_gpt}")
    print(f"原图背景: {use_original_bg}")
    print(f"元素遮罩: {mask_elements}")
    print()
    
    json_path, pptx_path = process_slide(
        image_path, output_dir,
        use_gpt=use_gpt,
        use_original_bg=use_original_bg,
        mask_elements=mask_elements
    )
    
    print(f"\n处理完成!")
    print(f"JSON: {json_path}")
    print(f"PPTX: {pptx_path}")
